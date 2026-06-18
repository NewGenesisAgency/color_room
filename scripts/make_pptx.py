#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Génère la présentation orale BTS CIEL IR — projet ColorRoom (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette ColorRoom ────────────────────────────────────────────────────────
INDIGO   = RGBColor(0x43, 0x61, 0xEE)
VIOLET   = RGBColor(0x7C, 0x3A, 0xED)
PINK     = RGBColor(0xEC, 0x48, 0x99)
DARK     = RGBColor(0x0F, 0x17, 0x2A)
SLATE    = RGBColor(0x33, 0x41, 0x55)
GREY     = RGBColor(0x64, 0x74, 0x8B)
LIGHT    = RGBColor(0xF1, 0xF4, 0xFB)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x06, 0xD6, 0xA0)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines = list of (text, size, color, bold, space_before)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold, sb) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if sb: p.space_before = Pt(sb)
        p.space_after = Pt(2)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt

def bullets(slide, x, y, w, h, items, size=18, color=SLATE, gap=10):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        # it = (text, level)  level 0/1
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = lvl
        bullet = "▸ " if lvl == 0 else "– "
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - (2 if lvl else 0))
        r.font.color.rgb = color if lvl == 0 else GREY
        r.font.bold = False; r.font.name = "Calibri"
    return tb

def header(slide, num, title, kicker):
    """Bandeau haut standard pour slides de contenu."""
    rect(slide, 0, 0, SW, Inches(1.25), DARK)
    rect(slide, 0, Inches(1.25), SW, Pt(4), VIOLET)
    # pastille numéro
    from pptx.enum.shapes import MSO_SHAPE
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(0.32), Inches(0.62), Inches(0.62))
    badge.fill.solid(); badge.fill.fore_color.rgb = VIOLET; badge.line.fill.background()
    badge.shadow.inherit = False
    btf = badge.text_frame; btf.word_wrap = False
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = str(num); br.font.size = Pt(22); br.font.bold = True; br.font.color.rgb = WHITE
    textbox(slide, Inches(1.25), Inches(0.20), Inches(11.5), Inches(0.95), [
        (kicker, 12, GREEN, True, 0),
        (title, 26, WHITE, True, 0),
    ])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, SW, SH, DARK)
# bande dégradée simulée (3 rectangles)
rect(s, 0, Inches(3.05), SW, Inches(0.10), INDIGO)
rect(s, 0, Inches(3.15), SW, Inches(0.10), VIOLET)
rect(s, 0, Inches(3.25), SW, Inches(0.10), PINK)
textbox(s, Inches(1), Inches(1.4), Inches(11.3), Inches(1.6), [
    ("PROJET BTS CIEL — OPTION INFORMATIQUE & RÉSEAUX", 14, GREEN, True, 0),
    ("ColorRoom", 60, WHITE, True, 6),
], align=PP_ALIGN.LEFT)
textbox(s, Inches(1), Inches(3.5), Inches(11.3), Inches(1.2), [
    ("Serious games interactifs sur 42 dalles LED pilotées en réseau", 22, RGBColor(0xCB,0xD5,0xE1), False, 0),
    ("Application web embarquée sur Raspberry Pi — mesure et pilotage de la lumière", 16, GREY, False, 6),
], align=PP_ALIGN.LEFT)
textbox(s, Inches(1), Inches(5.6), Inches(11.3), Inches(1.4), [
    ("Présenté par : [Ton Prénom NOM]", 16, WHITE, True, 0),
    ("Session 2026 — Lycée Édouard Branly", 14, GREY, False, 4),
    ("Partenaires : LUMEN Campus Lumière · ENTPE / LTDS", 13, GREY, False, 2),
], align=PP_ALIGN.LEFT)
notes(s, "Bonjour, je m'appelle [...], je suis en 2e année de BTS CIEL option Informatique et Réseaux. "
         "Je vais vous présenter mon projet : ColorRoom, une application qui transforme une salle de 42 dalles LED "
         "en plateforme de jeux éducatifs sur la couleur. Durée : 20 min de présentation, puis une démonstration, puis vos questions. "
         "ASTUCE : respire, parle lentement, regarde le jury.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Sommaire
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, Inches(0.25), SH, VIOLET)
textbox(s, Inches(0.8), Inches(0.5), Inches(11), Inches(1), [("Sommaire", 34, DARK, True, 0)])
col1 = ["1. Contexte & commanditaire","2. Problématique & objectifs","3. Cahier des charges & contraintes",
        "4. Architecture générale","5. Choix techniques justifiés","6. Réseau & déploiement (Docker / Pi)"]
col2 = ["7. Fonctionnalités du projet","8. Focus : éditeur, IA, mesure","9. Sécurité","10. Gestion de projet",
        "11. Difficultés & solutions","12. Bilan & perspectives"]
bullets(s, Inches(0.9), Inches(1.8), Inches(6), Inches(5), col1, size=20, gap=16)
bullets(s, Inches(7.0), Inches(1.8), Inches(6), Inches(5), col2, size=20, gap=16)
notes(s, "Voici le déroulé. Je commence par le contexte, puis l'architecture technique, le réseau et le déploiement, "
         "ensuite les fonctionnalités, la sécurité, la gestion de projet, et je termine par le bilan.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Contexte
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,1,"Contexte & commanditaire","LE PROJET")
bullets(s, Inches(0.7), Inches(1.7), Inches(7.4), Inches(5), [
    "La ColorRoom : une salle équipée de 42 dalles LED RVB pilotables individuellement",
    "Commanditaire : LUMEN Campus Lumière (Lyon), en lien avec l'ENTPE / laboratoire LTDS",
    "But : des « serious games » pour sensibiliser à la couleur, la lumière, la perception",
    ("Public : étudiants, visiteurs, démonstrations pédagogiques", 1),
    "Le matériel existe — il fallait le LOGICIEL pour l'animer et le rendre utile",
], size=18, gap=14)
card = rect(s, Inches(8.5), Inches(1.9), Inches(4.2), Inches(4.4), WHITE)
textbox(s, Inches(8.8), Inches(2.2), Inches(3.7), Inches(4), [
    ("EN CHIFFRES", 12, VIOLET, True, 0),
    ("42", 44, INDIGO, True, 8), ("dalles LED RVB", 14, GREY, False, 0),
    ("6 × 7", 30, VIOLET, True, 10), ("grille de la salle", 14, GREY, False, 0),
    ("1", 30, PINK, True, 10), ("Raspberry Pi (autonome)", 14, GREY, False, 0),
])
notes(s, "La ColorRoom est une installation physique : 42 dalles LED commandées une par une. "
         "Le commanditaire voulait des jeux éducatifs sur la couleur. Le matériel était là, mais sans logiciel "
         "pour le piloter de façon ludique : c'est tout l'objet de mon projet.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Problématique
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,2,"Problématique & objectifs","LE BESOIN")
# problématique en encart
rect(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.2), VIOLET)
textbox(s, Inches(1.0), Inches(1.85), Inches(11.4), Inches(1), [
    ("« Comment piloter et mesurer la lumière de 42 dalles LED, tout en proposant", 18, WHITE, True, 0),
    ("des jeux éducatifs créables facilement, sur un système autonome et hors-ligne ? »", 18, WHITE, True, 2),
])
bullets(s, Inches(0.7), Inches(3.3), Inches(11.9), Inches(3.6), [
    "Objectif 1 — Piloter les dalles en temps réel (couleur, intensité) via le réseau",
    "Objectif 2 — Mesurer réellement la lumière émise (colorimètre Konica CS-160)",
    "Objectif 3 — Offrir des jeux + un éditeur permettant d'en créer SANS coder",
    "Objectif 4 — Fonctionner 100 % en local sur Raspberry Pi (pas de cloud obligatoire)",
    "Objectif 5 — Multijoueur (chaque joueur pilote une dalle depuis son téléphone)",
], size=18, gap=15)
notes(s, "La problématique tient en une phrase : piloter ET mesurer la lumière, avec des jeux créables facilement, "
         "sur un système autonome. J'en ai tiré 5 objectifs concrets que je vais détailler.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Cahier des charges
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,3,"Cahier des charges & contraintes","CADRE TECHNIQUE")
# deux colonnes : fonctionnel / contraintes
textbox(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.5),[("Exigences fonctionnelles",18,VIOLET,True,0)])
bullets(s, Inches(0.7), Inches(2.15), Inches(6), Inches(4.6), [
    "Jeux jouables sur les dalles + tablette",
    "Éditeur de jeux low-code (blocs) + code Python/JS",
    "Mode multijoueur en réseau local",
    "Mesure colorimétrique et diagramme CIE",
    "Interface d'administration (comptes, classes)",
], size=16, gap=12)
textbox(s, Inches(6.9), Inches(1.6), Inches(6), Inches(0.5),[("Contraintes techniques",18,PINK,True,0)])
bullets(s, Inches(6.9), Inches(2.15), Inches(5.9), Inches(4.6), [
    "Cible : Raspberry Pi (ressources limitées)",
    "Fonctionnement hors-ligne possible",
    "API « Supervision » imposée pour les dalles",
    "Colorimètre Konica Minolta CS-160 (SDK)",
    "Déploiement reproductible et maintenable",
], size=16, gap=12)
notes(s, "Côté fonctionnel : les jeux, l'éditeur, le multijoueur, la mesure, l'admin. "
         "Côté contraintes : tout doit tenir sur un Raspberry Pi, marcher hors-ligne, "
         "utiliser l'API Supervision imposée pour les dalles et le SDK du colorimètre CS-160.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,WHITE); header(s,4,"Architecture générale","VUE D'ENSEMBLE")
from pptx.enum.shapes import MSO_SHAPE
def box(x,y,w,h,title,sub,color):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit=False
    tf = sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=title; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=WHITE
    if sub:
        p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
        r2=p2.add_run(); r2.text=sub; r2.font.size=Pt(10); r2.font.color.rgb=RGBColor(0xE2,0xE8,0xF0)
    return sp
def arrow(x,y,w,h):
    a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x,y,w,h)
    a.fill.solid(); a.fill.fore_color.rgb=GREY; a.line.fill.background(); a.shadow.inherit=False
Y=Inches(2.4)
box(Inches(0.6),Y,Inches(2.6),Inches(1.4),"Tablette / PC","Navigateur (UI)",INDIGO)
arrow(Inches(3.3),Inches(2.85),Inches(0.7),Inches(0.5))
box(Inches(4.1),Y,Inches(3.0),Inches(1.4),"App ColorRoom","Next.js · Docker · Pi",VIOLET)
arrow(Inches(7.2),Inches(2.85),Inches(0.7),Inches(0.5))
box(Inches(8.0),Y,Inches(2.4),Inches(1.4),"API Supervision","HTTP local",SLATE)
arrow(Inches(10.5),Inches(2.85),Inches(0.7),Inches(0.5))
box(Inches(11.3),Y,Inches(1.5),Inches(1.4),"42 dalles","LED RVB",GREEN)
# couche basse : IA + CS-160 + DB
box(Inches(4.1),Inches(4.4),Inches(3.0),Inches(1.1),"IA","Gemini ↔ Ollama",PINK)
box(Inches(0.6),Inches(4.4),Inches(2.6),Inches(1.1),"SQLite","jeux · scores · users",AMBER)
box(Inches(8.0),Inches(4.4),Inches(2.4),Inches(1.1),"CS-160","colorimètre (USB)",SLATE)
textbox(s, Inches(0.6), Inches(5.8), Inches(12), Inches(1), [
    ("Tout est conteneurisé (Docker Compose) et tourne en local sur le Raspberry Pi — l'IA cloud (Gemini) n'est qu'une option, le local (Ollama) prend le relais hors-ligne.", 14, GREY, False, 0)
])
notes(s, "Voici le schéma global. L'utilisateur ouvre l'app dans un navigateur. L'app, en Next.js et conteneurisée, "
         "tourne sur le Pi. Pour allumer les dalles, elle parle à l'API Supervision en HTTP. En parallèle : "
         "une base SQLite pour les données, le colorimètre CS-160 en USB, et l'IA (cloud Gemini ou local Ollama) pour générer des jeux.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Choix techniques
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,5,"Choix techniques justifiés","POURQUOI CES OUTILS")
rows = [
    ("Next.js / React + TypeScript","Front + back dans un seul projet, typage strict = moins de bugs"),
    ("SQLite (better-sqlite3)","Base embarquée, zéro serveur à installer — parfait sur Pi"),
    ("Docker Compose","Déploiement reproductible : 1 commande lance toute la stack"),
    ("Three.js","Vue 3D de la salle en temps réel dans le navigateur"),
    ("Ollama (IA locale)","Génération de jeux 100 % hors-ligne, données privées"),
]
y = Inches(1.7)
for i,(tech,why) in enumerate(rows):
    rect(s, Inches(0.7), y, Inches(11.9), Inches(0.86), WHITE if i%2==0 else RGBColor(0xE9,0xED,0xF7))
    textbox(s, Inches(0.9), y+Emu(40000), Inches(4.2), Inches(0.8), [(tech,16,VIOLET,True,0)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(5.2), y+Emu(40000), Inches(7.2), Inches(0.8), [(why,15,SLATE,False,0)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.92)
notes(s, "Chaque choix est justifié. Next.js avec TypeScript pour un code typé et un seul projet front+back. "
         "SQLite parce qu'elle est embarquée, idéale sur Pi. Docker pour un déploiement reproductible en une commande. "
         "Three.js pour la 3D. Et Ollama pour faire tourner une IA en local, sans cloud.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7bis — Pourquoi React/Next/Node/TS vs JS pur + Node-RED
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT)
header(s,5,"Pourquoi React / Next.js / Node.js / TypeScript ?","JUSTIFICATION DU CHOIX (vs JS pur + Node-RED)")
# Colonne gauche : approche initiale (limites)
rect(s, Inches(0.7), Inches(1.65), Inches(5.85), Inches(0.6), SLATE)
textbox(s, Inches(0.9), Inches(1.72), Inches(5.5), Inches(0.5), [("Approche initiale : Node-RED + JS pur",16,WHITE,True,0)])
bullets(s, Inches(0.8), Inches(2.4), Inches(5.7), Inches(4.4), [
    "Node-RED : génial pour prototyper l'IoT, mais limité pour une vraie UI riche (3D, éditeur de jeux)",
    "Flux stockés en JSON → difficiles à versionner et à relire dans Git",
    "JS pur : aucun typage → erreurs détectées seulement à l'exécution",
    "Refactorisation risquée, autocomplétion faible sur un gros projet",
    "Montée en complexité vite ingérable (logique métier importante)",
], size=14, color=SLATE, gap=11)
# Colonne droite : choix retenu (avantages)
rect(s, Inches(6.75), Inches(1.65), Inches(5.85), Inches(0.6), VIOLET)
textbox(s, Inches(6.95), Inches(1.72), Inches(5.5), Inches(0.5), [("Choix retenu : Next.js + React + Node + TS",16,WHITE,True,0)])
bullets(s, Inches(6.85), Inches(2.4), Inches(5.7), Inches(4.4), [
    "React : interface en composants réutilisables, idéale pour une UI complexe",
    "Next.js : front + back (routes API) dans UN seul projet, routing intégré",
    "Node.js : même langage côté serveur et navigateur (un seul écosystème)",
    "TypeScript : typage strict → bugs attrapés à la COMPILATION, refacto sûre",
    "Tout est du code texte → versionnable proprement, testable (tsc + build)",
], size=14, color=SLATE, gap=11)
notes(s, "Question classique du jury : pourquoi ne pas être resté sur du Node-RED et du JavaScript pur ? "
         "Node-RED est excellent pour prototyper de l'IoT par flux, mais dès qu'on veut une vraie interface riche — "
         "la vue 3D, un éditeur de jeux, du multijoueur — il montre ses limites, et ses flux en JSON sont durs à versionner. "
         "Le JavaScript pur, lui, n'a pas de typage : les erreurs n'apparaissent qu'à l'exécution. "
         "J'ai donc choisi React pour des composants réutilisables, Next.js pour réunir front et back dans un seul projet, "
         "Node.js pour un seul langage partout, et surtout TypeScript : le typage attrape les bugs dès la compilation "
         "et rend la refactorisation sûre. Le tout est du code texte, donc versionnable et testable proprement.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Réseau & déploiement
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,6,"Réseau & déploiement","DOCKER / RASPBERRY PI")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.4), Inches(5), [
    "Orchestration via docker-compose.yml (4 services)",
    ("color-room : l'application (port 8080)", 1),
    ("ollama + ollama-pull : IA locale (port 11434)", 1),
    ("portainer : supervision des conteneurs (9000)", 1),
    "Réseau local : tablettes/téléphones sur le même Wi-Fi",
    "Variables sensibles isolées dans un fichier .env (non versionné)",
    "Mise à jour : git pull + docker compose up -d --build",
], size=16, gap=11)
card = rect(s, Inches(7.5), Inches(1.9), Inches(5.1), Inches(4.2), DARK)
textbox(s, Inches(7.8), Inches(2.15), Inches(4.6), Inches(3.8), [
    ("# déploiement sur le Pi", 13, GREEN, True, 0),
    ("git pull origin main", 14, WHITE, False, 8),
    ("docker compose up -d --build", 14, WHITE, False, 4),
    ("", 8, WHITE, False, 4),
    ("# 4 conteneurs démarrent :", 13, GREEN, True, 6),
    ("• color-room   :8080", 13, RGBColor(0xCB,0xD5,0xE1), False, 4),
    ("• ollama       :11434", 13, RGBColor(0xCB,0xD5,0xE1), False, 2),
    ("• portainer    :9000", 13, RGBColor(0xCB,0xD5,0xE1), False, 2),
])
notes(s, "Le déploiement est entièrement géré par Docker Compose : 4 services. L'application sur le port 8080, "
         "l'IA locale Ollama, et Portainer pour superviser les conteneurs. Les téléphones rejoignent via le Wi-Fi local. "
         "Les secrets (clé API, mot de passe admin) sont dans un fichier .env qui n'est jamais versionné sur Git. "
         "Mettre à jour le Pi = git pull puis docker compose up.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Fonctionnalités
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,WHITE); header(s,7,"Fonctionnalités du projet","CE QUE FAIT L'APPLI")
feats = [
    ("🎮","Jeux intégrés","Color Speed, Simon, Puissance 4, Mémoire, mesure…",INDIGO),
    ("🧩","Éditeur low-code","Créer un jeu en reliant des blocs, ou en Python/JS",VIOLET),
    ("🤖","Génération par IA","Décrire un jeu en une phrase → il est créé",PINK),
    ("📱","Multijoueur","1 joueur = 1 dalle, rejoint via QR code",GREEN),
    ("🎨","Mesure CS-160","Colorimètre réel + diagramme CIE 1931",AMBER),
    ("🔐","Administration","Comptes, classes, scores, niveaux",SLATE),
]
x0,y0=Inches(0.7),Inches(1.7); cw,ch=Inches(3.9),Inches(2.3)
for i,(emo,t,d,c) in enumerate(feats):
    cx = x0 + (i%3)*Inches(4.05)
    cy = y0 + (i//3)*Inches(2.5)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch)
    card.fill.solid(); card.fill.fore_color.rgb=LIGHT; card.line.color.rgb=c; card.line.width=Pt(1.5); card.shadow.inherit=False
    textbox(s, cx+Inches(0.25), cy+Inches(0.2), cw-Inches(0.5), ch-Inches(0.4), [
        (emo+"  "+t, 17, c, True, 0),
        (d, 13, SLATE, False, 8),
    ])
notes(s, "Six grandes fonctionnalités. Des jeux prêts à jouer. Un éditeur low-code pour en créer sans coder, "
         "ou en Python/JS pour aller plus loin. La génération de jeux par IA. Le multijoueur où chacun pilote une dalle "
         "avec son téléphone. La mesure réelle au colorimètre avec diagramme CIE. Et l'administration des comptes.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Focus éditeur
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,8,"Focus technique — L'éditeur de jeux","POINT FORT 1/3")
bullets(s, Inches(0.7), Inches(1.7), Inches(7.2), Inches(5), [
    "Programmation visuelle par NŒUDS reliés (« blueprint »)",
    ("Évènements (clic dalle, timer…) → actions (allumer, score, son)", 1),
    "Catalogue de ~190 blocs : logique, maths, rendu, animations, multijoueur",
    "Bloc « Script » : code Python OU JavaScript exécuté dans le jeu",
    "Aperçu en direct sur la vue 3D + envoi aux vraies dalles",
    "Les jeux créés sont sauvegardés en base et jouables par tous",
], size=16, gap=12)
card = rect(s, Inches(8.2), Inches(1.9), Inches(4.4), Inches(4.3), WHITE)
textbox(s, Inches(8.5), Inches(2.15), Inches(3.9), Inches(4), [
    ("EXEMPLE DE FLUX", 12, VIOLET, True, 0),
    ("● Démarrer", 15, INDIGO, True, 10),
    ("   ↓", 13, GREY, False, 2),
    ("● Dalle aléatoire", 15, VIOLET, True, 2),
    ("   ↓", 13, GREY, False, 2),
    ("● Clic joueur ?", 15, PINK, True, 2),
    ("   ↓", 13, GREY, False, 2),
    ("● +1 point + son", 15, GREEN, True, 2),
])
notes(s, "Le cœur technique : l'éditeur. C'est de la programmation visuelle par blocs reliés, comme dans Unreal Engine. "
         "On part d'un évènement (un clic sur une dalle) et on enchaîne des actions. Il y a près de 190 blocs. "
         "Pour les utilisateurs avancés, un bloc Script permet d'écrire du Python ou du JavaScript. "
         "Tout s'affiche en direct sur la 3D et sur les vraies dalles.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Focus IA
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,9,"Focus technique — Génération par IA","POINT FORT 2/3")
bullets(s, Inches(0.7), Inches(1.7), Inches(7.2), Inches(5), [
    "L'utilisateur décrit le jeu en langage naturel",
    "Cascade intelligente : Gemini (cloud) → repli Ollama (local)",
    ("Marche hors-ligne grâce au modèle local sur le Pi", 1),
    "L'IA renvoie un JSON strict → converti en blocs + interface",
    "Garde-fou : refuse les jeux impossibles (Forza, GTA…) avec un message clair",
    "Validation/sanitisation systématique de la sortie (sécurité + robustesse)",
], size=16, gap=12)
card = rect(s, Inches(8.2), Inches(1.9), Inches(4.4), Inches(4.3), DARK)
textbox(s, Inches(8.5), Inches(2.15), Inches(3.9), Inches(4), [
    ("CASCADE IA", 12, GREEN, True, 0),
    ("1. Clé Gemini ?", 14, WHITE, True, 10),
    ("   → oui : Gemini Flash", 12, RGBColor(0xCB,0xD5,0xE1), False, 2),
    ("2. Échec / hors-ligne ?", 14, WHITE, True, 8),
    ("   → Ollama (Pi)", 12, RGBColor(0xCB,0xD5,0xE1), False, 2),
    ("3. Sortie JSON validée", 14, WHITE, True, 8),
    ("   → blocs + UI générés", 12, RGBColor(0xCB,0xD5,0xE1), False, 2),
])
notes(s, "Deuxième point fort : l'IA. On décrit un jeu en une phrase, et l'IA le construit. "
         "J'ai mis en place une cascade : d'abord Gemini en cloud si une clé existe, sinon repli automatique sur Ollama "
         "qui tourne en local sur le Pi — donc ça marche même hors-ligne. La sortie de l'IA est du JSON strict, "
         "que je valide et nettoie avant de l'utiliser. Et un garde-fou refuse poliment les jeux irréalisables.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Focus mesure CS-160
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,10,"Focus technique — Mesure colorimétrique","POINT FORT 3/3")
bullets(s, Inches(0.7), Inches(1.7), Inches(7.2), Inches(5), [
    "Colorimètre Konica Minolta CS-160 piloté via une API dédiée",
    "Mesure réelle de la lumière émise par les dalles (X, Y, Z / x, y, Lv)",
    "Affichage sur le diagramme de chromaticité CIE 1931",
    "Jeux pédagogiques basés sur la mesure (reconnaître une couleur, métamérie)",
    "Calcul d'écart de couleur (ΔE) pour scorer la précision du joueur",
    "Lien direct entre le virtuel (jeu) et le réel (physique mesurée)",
], size=16, gap=12)
card = rect(s, Inches(8.2), Inches(1.9), Inches(4.4), Inches(4.3), WHITE)
textbox(s, Inches(8.5), Inches(2.15), Inches(3.9), Inches(4), [
    ("CHAÎNE DE MESURE", 12, VIOLET, True, 0),
    ("Dalle allumée", 15, INDIGO, True, 12),
    ("   ↓ émet de la lumière", 12, GREY, False, 2),
    ("CS-160 mesure", 15, VIOLET, True, 8),
    ("   ↓ x, y, Lv", 12, GREY, False, 2),
    ("Diagramme CIE", 15, PINK, True, 8),
    ("   ↓ ΔE", 12, GREY, False, 2),
    ("Score du joueur", 15, GREEN, True, 8),
])
notes(s, "Troisième point fort, très « métier lumière » : la mesure. Le colorimètre CS-160 mesure réellement la lumière "
         "émise par les dalles. On affiche le résultat sur le diagramme CIE 1931, le standard de la colorimétrie. "
         "Certains jeux notent le joueur selon l'écart de couleur ΔE. C'est ce qui relie le virtuel au réel.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Sécurité
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,11,"Sécurité","BONNES PRATIQUES")
bullets(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(5), [
    "Mots de passe hachés (PBKDF2-SHA512, 100 000 itérations, sel aléatoire) — jamais en clair",
    "Sessions par cookie httpOnly (30 jours), rôles utilisateur (élève / enseignant / admin)",
    "Secrets (clé API, mot de passe admin) dans .env — exclu de Git via .gitignore",
    "Transactions SQL atomiques (ACID) pour éviter les états incohérents",
    "Réseau local isolé : pas d'exposition sur Internet, l'IA peut tourner 100 % hors-ligne",
    "RETOUR D'EXPÉRIENCE : un secret committé par erreur → purge de l'historique Git (git filter-repo) + rotation",
], size=16, gap=13)
notes(s, "La sécurité, important en CIEL. Les mots de passe sont hachés avec PBKDF2, jamais stockés en clair. "
         "Sessions par cookie sécurisé, gestion de rôles. Les secrets sont dans un .env hors Git. "
         "Les écritures critiques en base sont des transactions atomiques. Et un vrai retour d'expérience : "
         "j'ai un jour committé un secret par erreur, je l'ai purgé de tout l'historique Git avec filter-repo. "
         "Ça montre que je connais les bonnes pratiques de gestion des secrets.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Gestion de projet
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,12,"Gestion de projet","MÉTHODE & OUTILS")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.2), Inches(5), [
    "Versioning Git + GitHub (branches de travail, commits clairs)",
    "Développement incrémental, fonctionnalité par fonctionnalité",
    "Tests à chaque étape : typage (tsc) + build de production",
    "Documentation : 15 diagrammes UML + guide technique",
    "Conteneurisation pour des déploiements fiables et répétables",
], size=16, gap=13)
card = rect(s, Inches(7.2), Inches(1.9), Inches(5.4), Inches(4.2), WHITE)
textbox(s, Inches(7.5), Inches(2.15), Inches(4.9), Inches(3.9), [
    ("OUTILS UTILISÉS", 12, VIOLET, True, 0),
    ("Code      ·  VS Code / TypeScript", 14, SLATE, False, 10),
    ("Versions  ·  Git + GitHub", 14, SLATE, False, 6),
    ("Conteneurs·  Docker Compose", 14, SLATE, False, 6),
    ("Supervision· Portainer", 14, SLATE, False, 6),
    ("Modélisation· UML (PlantUML)", 14, SLATE, False, 6),
    ("Cible     ·  Raspberry Pi", 14, SLATE, False, 6),
])
notes(s, "Côté gestion : tout est versionné sur Git et GitHub avec des branches et des commits clairs. "
         "J'ai développé de façon incrémentale, en testant à chaque étape (vérification de types et build complet). "
         "J'ai documenté avec 15 diagrammes UML et un guide technique. Et tout est conteneurisé pour des déploiements fiables.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Difficultés
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,WHITE); header(s,13,"Difficultés rencontrées & solutions","DÉMARCHE D'INGÉNIEUR")
probs = [
    ("Lag de la vue 3D sur tablette","Textures lourdes supprimées, pixel ratio plafonné, snapshots dédupliqués"),
    ("Contextes WebGL bloqués (crash)","Libération explicite du contexte (forceContextLoss) + repli gracieux"),
    ("Son d'un jeu joué « à l'infini »","Bug d'unité ms/secondes corrigé dans la génération audio"),
    ("Tetris injouable sur 6×7","Refonte en jeu de combinaison de couleurs, pièces de 1-2 cases"),
]
y=Inches(1.75)
for i,(p,sol) in enumerate(probs):
    rect(s, Inches(0.7), y, Inches(5.7), Inches(1.0), RGBColor(0xFE,0xE2,0xE2))
    textbox(s, Inches(0.9), y+Emu(30000), Inches(5.3), Inches(0.9), [("⚠ "+p,14,RGBColor(0xB9,0x1C,0x1C),True,0)], anchor=MSO_ANCHOR.MIDDLE)
    arrow_x=Inches(6.45)
    a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,arrow_x,y+Inches(0.32),Inches(0.4),Inches(0.36))
    a.fill.solid(); a.fill.fore_color.rgb=GREY; a.line.fill.background(); a.shadow.inherit=False
    rect(s, Inches(6.95), y, Inches(5.65), Inches(1.0), RGBColor(0xDC,0xFC,0xE7))
    textbox(s, Inches(7.15), y+Emu(30000), Inches(5.3), Inches(0.9), [("✓ "+sol,13,RGBColor(0x15,0x80,0x3D),True,0)], anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(1.15)
notes(s, "Un projet, c'est surtout des problèmes à résoudre. La 3D ramait sur tablette : j'ai allégé les textures et "
         "réduit le travail de rendu. Des crashs WebGL : je libère maintenant explicitement le contexte graphique. "
         "Un son qui durait à l'infini : une confusion millisecondes/secondes que j'ai corrigée. "
         "Et le Tetris injouable sur une si petite grille : je l'ai transformé en jeu de combinaison de couleurs. "
         "C'est ça, la démarche d'ingénieur : diagnostiquer puis corriger.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Bilan & perspectives
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,LIGHT); header(s,14,"Bilan & perspectives","CONCLUSION TECHNIQUE")
textbox(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.5),[("✓ Objectifs atteints",18,GREEN,True,0)])
bullets(s, Inches(0.7), Inches(2.15), Inches(6), Inches(4.5), [
    "Pilotage temps réel des 42 dalles",
    "Mesure réelle au colorimètre + CIE",
    "Éditeur + génération IA fonctionnels",
    "Multijoueur en réseau local",
    "Déploiement automatisé sur Pi",
], size=16, gap=11)
textbox(s, Inches(6.9), Inches(1.6), Inches(6), Inches(0.5),[("→ Perspectives",18,VIOLET,True,0)])
bullets(s, Inches(6.9), Inches(2.15), Inches(5.9), Inches(4.5), [
    "Vrai système de salons multijoueurs (codes)",
    "Streaming de la génération IA (temps réel)",
    "Plus de jeux pédagogiques et de tutoriels",
    "Tableau de bord enseignant enrichi",
    "Tests automatisés (CI/CD)",
], size=16, gap=11)
notes(s, "Bilan : les 5 objectifs sont atteints — pilotage, mesure, éditeur + IA, multijoueur, déploiement. "
         "Pour la suite : un vrai système de salons avec codes, le streaming de la génération IA en temps réel, "
         "plus de jeux, un tableau de bord enseignant, et de l'intégration continue pour automatiser les tests.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Démonstration
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,DARK)
rect(s, 0, Inches(3.4), SW, Pt(4), VIOLET)
textbox(s, Inches(1), Inches(2.4), Inches(11.3), Inches(2), [
    ("PLACE À LA", 16, GREEN, True, 0),
    ("Démonstration", 50, WHITE, True, 6),
], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.0), Inches(11.3), Inches(2.5), [
    ("1. Lancer un jeu sur les dalles  ·  2. Créer un jeu (éditeur)", 16, RGBColor(0xCB,0xD5,0xE1), False, 0),
    ("3. Générer un jeu avec l'IA  ·  4. Multijoueur (QR + téléphone)", 16, RGBColor(0xCB,0xD5,0xE1), False, 6),
    ("5. Mesure CS-160 + diagramme CIE  ·  6. Portainer / Docker", 16, RGBColor(0xCB,0xD5,0xE1), False, 6),
], align=PP_ALIGN.CENTER)
notes(s, "DÉMO — ordre conseillé : 1) lancer Color Speed sur les dalles, 2) créer un petit jeu dans l'éditeur, "
         "3) en générer un avec l'IA, 4) montrer le multijoueur avec le QR code et un téléphone, "
         "5) faire une mesure au CS-160 et la voir sur le diagramme CIE, 6) finir sur Portainer pour montrer les conteneurs. "
         "PRÉPARE UN PLAN B : capture vidéo/écran au cas où le matériel bug le jour J.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Merci / Questions
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); rect(s,0,0,SW,SH,DARK)
rect(s, 0, Inches(2.0), SW, Inches(0.08), INDIGO)
rect(s, 0, Inches(2.08), SW, Inches(0.08), VIOLET)
rect(s, 0, Inches(2.16), SW, Inches(0.08), PINK)
textbox(s, Inches(1), Inches(2.8), Inches(11.3), Inches(2), [
    ("Merci de votre attention", 44, WHITE, True, 0),
    ("Avez-vous des questions ?", 22, GREEN, True, 14),
], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.6), [
    ("ColorRoom — BTS CIEL option IR — Session 2026", 13, GREY, False, 0),
], align=PP_ALIGN.CENTER)
notes(s, "Merci. Je suis à votre disposition pour vos questions. "
         "Garde en tête les questions probables : pourquoi Next.js ? comment sécurises-tu le réseau ? "
         "comment ça marche hors-ligne ? que se passe-t-il si une dalle tombe en panne ? combien d'utilisateurs simultanés ?")

prs.save("ColorRoom_Presentation_BTS_CIEL.pptx")
print("OK — fichier généré.")
